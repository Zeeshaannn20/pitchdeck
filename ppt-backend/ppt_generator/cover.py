from __future__ import annotations

import os
import base64
import logging
import tempfile
from typing import Optional

from PIL import Image
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

from .config import (
    CONTENT_FONT_NAME, TITLE_COLOR, GOLD_ACCENT,
)
from .utils import _emu_to_in
from .template import (
    get_blank_layout, remove_placeholders, apply_institute_slide_background,
)

logger = logging.getLogger(__name__)

def _save_base64_image(b64_str: str, temp_files: list[str],
                       suffix: str = ".png") -> Optional[str]:
    if not b64_str:
        return None
    try:
        img_bytes = base64.b64decode(b64_str)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        with open(tmp.name, "wb") as f:
            f.write(img_bytes)
        temp_files.append(tmp.name)
        return tmp.name
    except Exception as e:
        logger.warning(f"Could not decode instructor image: {e}")
        return None

def create_ioi_cover_slide(prs, subject_text: str, chapter_text: str,
                            instructor_name: str = "", qualification: str = "",
                            instructor_image_b64: Optional[str] = None,
                            logo_path: str = "logo.png", batch_code: str = "",
                            temp_files: Optional[list] = None) -> None:
    if temp_files is None:
        temp_files = []

    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_placeholders(slide)
    apply_institute_slide_background(slide, prs=prs)

    sw = _emu_to_in(prs.slide_width)
    sh = _emu_to_in(prs.slide_height)

    logo_top = 0.50
    target_w = 2.60
    target_h = 0.80
    try:
        if os.path.exists(logo_path):
            with Image.open(logo_path) as im:
                w, h = im.size
            aspect   = (h / float(w)) if w else 1.0
            target_h = min(0.85, target_w * aspect)
            target_w = target_h / aspect
            slide.shapes.add_picture(
                logo_path,
                Inches((sw - target_w) / 2.0), Inches(logo_top),
                width=Inches(target_w), height=Inches(target_h),
            )
    except Exception as e:
        logger.warning(f"Cover logo failed: {e}")

    pill_w, pill_h = 3.80, 0.48
    pill_top  = logo_top + target_h + 0.25
    pill_left = (sw - pill_w) / 2.0
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(pill_left), Inches(pill_top),
                                  Inches(pill_w), Inches(pill_h))
    pill.adjustments[0] = 0.30
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
    pill.line.color.rgb = GOLD_ACCENT; pill.line.width = Pt(1.6)
    ptf = pill.text_frame; ptf.clear()
    ptf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p   = ptf.paragraphs[0]; p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    r   = p.add_run(); r.text = f"Batch Code: {batch_code}".strip()
    r.font.name = CONTENT_FONT_NAME; r.font.size = Pt(17)
    r.font.bold = True; r.font.color.rgb = TITLE_COLOR

    pic_left, pic_top, pic_w = 0.60, 2.40, 4.80
    ph_in    = 0.0
    img_path = (_save_base64_image(instructor_image_b64, temp_files)
                if instructor_image_b64 else None)
    if img_path and os.path.exists(img_path):
        try:
            with Image.open(img_path) as img:
                iw, ih = img.size
            aspect = ih / iw
            pic_h  = pic_w * aspect
            slide.shapes.add_picture(img_path, Inches(pic_left), Inches(pic_top),
                                     width=Inches(pic_w), height=Inches(pic_h))
            ph_in  = pic_h
            frame  = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(pic_left-0.10), Inches(pic_top-0.10),
                                            Inches(pic_w+0.20), Inches(pic_h+0.20))
            frame.adjustments[0] = 0.08
            frame.fill.background()
            frame.line.color.rgb = RGBColor(255, 255, 255); frame.line.width = Pt(2)
            slide.shapes._spTree.remove(frame.element)
            slide.shapes._spTree.insert(0, frame.element)
        except Exception as e:
            logger.warning(f"Instructor image failed: {e}")

    card_left, card_top, card_w, card_h = 5.50, 3.90, 7.05, 2.15
    info = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(card_left), Inches(card_top),
                                  Inches(card_w), Inches(card_h))
    info.adjustments[0] = 0.10
    info.fill.background()
    info.line.color.rgb = RGBColor(255, 255, 255); info.line.width = Pt(3)
    itf = info.text_frame; itf.clear()
    itf.margin_left = Inches(0.55); itf.margin_right  = Inches(0.55)
    itf.margin_top  = Inches(0.26); itf.margin_bottom = Inches(0.20)
    itf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    def _row(label, value, first=False):
        pr = (itf.paragraphs[0]
              if (first and not itf.paragraphs[0].text)
              else itf.add_paragraph())
        pr.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT; pr.space_after = Pt(12)
        rl = pr.add_run(); rl.text = f"{label}: "
        rl.font.name = CONTENT_FONT_NAME; rl.font.size = Pt(21)
        rl.font.bold = True; rl.font.color.rgb = RGBColor(255, 255, 255)
        rv = pr.add_run(); rv.text = value or ""
        rv.font.name = CONTENT_FONT_NAME; rv.font.size = Pt(21)
        rv.font.bold = False; rv.font.color.rgb = RGBColor(255, 255, 255)

    _row("Subject", subject_text or "", first=True)
    _row("Chapter", chapter_text or "")

    name_top  = (pic_top + ph_in + 0.15) if ph_in else 5.50
    name_left, name_w, name_h = pic_left, 4.80, 1.28
    nb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(name_left), Inches(name_top),
                                Inches(name_w), Inches(name_h))
    nb.adjustments[0] = 0.14
    nb.fill.solid(); nb.fill.fore_color.rgb = RGBColor(255, 255, 255)
    nb.line.color.rgb = GOLD_ACCENT; nb.line.width = Pt(2)
    nft = nb.text_frame; nft.clear()
    nft.margin_left = Inches(0.40); nft.margin_right  = Inches(0.40)
    nft.margin_top  = Inches(0.10); nft.margin_bottom = Inches(0.06)
    nft.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    pn = nft.paragraphs[0]
    pn.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT; pn.space_after = Pt(2)
    rn = pn.add_run(); rn.text = (instructor_name or "Instructor").strip()
    rn.font.name = CONTENT_FONT_NAME; rn.font.size = Pt(32)
    rn.font.bold = True; rn.font.color.rgb = TITLE_COLOR

    pq = nft.add_paragraph(); pq.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    rq = pq.add_run(); rq.text = (qualification or "").strip()
    rq.font.name = CONTENT_FONT_NAME; rq.font.size = Pt(16)
    rq.font.bold = False; rq.font.color.rgb = RGBColor(0, 0, 0)
