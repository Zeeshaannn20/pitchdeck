from __future__ import annotations

import os
import re
import logging
import tempfile
from typing import Optional

from PIL import Image, ImageFont, ImageDraw
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

from .config import (
    MARGIN_LEFT, CONTENT_FONT_NAME, TITLE_FONT_NAME, MONO_FONT,
    CONTENT_FONT_SIZE, TITLE_FONT_SIZE,
    TITLE_COLOR, CONTENT_COLOR, ACCENT_PRIMARY, ACCENT_SECONDARY,
    WARNING_COLOR, ERROR_COLOR, CODE_BG, CODE_BORDER, CODE_TEXT,
    COLOR_BLACK, GOLD_ACCENT, SLIDE_BG_PRIMARY,
    BODY_LINE_SPACING_PT, CODE_LINE_SPACING_PT,
    PARA_SPACE_BEFORE_PT, PARA_SPACE_AFTER_PT,
)
from .utils import (
    _emu_to_in, _measure_text_inches, _estimate_paragraph_height_in,
    _content_width_in, _slide_box, SlideCtx,
)
from .template import (
    get_blank_layout, remove_placeholders, apply_institute_slide_background,
    add_institute_logo, create_institute_title_area, setup_text_frame_margins,
)
from .parsers import (
    normalize_separators, extract_code_block, parse_table_from_text,
    extract_inline_pipe_tables,
)

logger = logging.getLogger(__name__)

# =============================
# Paragraph styling helpers
# =============================
def style_paragraph(paragraph, line_pt=BODY_LINE_SPACING_PT,
                    before_pt=PARA_SPACE_BEFORE_PT, after_pt=PARA_SPACE_AFTER_PT,
                    align=None) -> None:
    try: paragraph.line_spacing = Pt(line_pt)
    except Exception: pass
    try:
        paragraph.space_before = Pt(before_pt)
        paragraph.space_after  = Pt(after_pt)
    except Exception: pass
    if align is not None:
        try: paragraph.alignment = align
        except Exception: pass

def _label_run(para, label: str, color: RGBColor) -> None:
    r = para.add_run()
    r.text = f"{label}: "
    r.font.name = CONTENT_FONT_NAME; r.font.size = CONTENT_FONT_SIZE
    r.font.bold = True; r.font.color.rgb = color

def _content_run(para, text: str) -> None:
    r = para.add_run()
    r.text = text
    r.font.name = CONTENT_FONT_NAME; r.font.size = CONTENT_FONT_SIZE
    r.font.color.rgb = CONTENT_COLOR

def _split_after_colon(text: str) -> str:
    return (text.split(":", 1)[1].strip() if ":" in text else text).strip()

# ============================
# Semantic item renderers
# ============================
def render_definition_item(para, text: str) -> None:
    style_paragraph(para)
    _label_run(para, "Definition", GOLD_ACCENT)
    _content_run(para, _split_after_colon(text))

def render_example_item(para, text: str) -> None:
    style_paragraph(para)
    _label_run(para, "Example", ACCENT_SECONDARY)
    _content_run(para, _split_after_colon(text))

def render_note_item(para, text: str) -> None:
    style_paragraph(para)
    _label_run(para, "Note", WARNING_COLOR)
    _content_run(para, _split_after_colon(text))

def render_warning_item(para, text: str) -> None:
    style_paragraph(para)
    label = "Warning" if text.lower().startswith("warning") else "Common Mistake"
    _label_run(para, label, ERROR_COLOR)
    _content_run(para, _split_after_colon(text))

def render_standard_item(para, text: str) -> None:
    style_paragraph(para)
    for part in re.split(r"(\*\*[^*]+\*\*)", text or ""):
        r = para.add_run()
        if part.startswith("**") and part.endswith("**"):
            r.text = part[2:-2]; r.font.bold = True; r.font.color.rgb = GOLD_ACCENT
        else:
            r.text = part; r.font.color.rgb = CONTENT_COLOR
        r.font.name = CONTENT_FONT_NAME; r.font.size = CONTENT_FONT_SIZE

_SEMANTIC_PREFIXES = (
    "definition:", "example:", "note:", "warning:", "common mistake:",
    "real-world use case:", "use case:", "practical tip:", "exam tip:",
)

def render_content_item(para, text: str) -> None:
    if not text:
        return
    low = text.lower().strip()
    if low.startswith("definition:"):
        render_definition_item(para, text)
    elif low.startswith("example:"):
        render_example_item(para, text)
    elif low.startswith("note:"):
        render_note_item(para, text)
    elif low.startswith(("warning:", "common mistake:")):
        render_warning_item(para, text)
    elif low.startswith(("real-world use case:", "use case:")):
        style_paragraph(para)
        _label_run(para, "Use Case", ACCENT_SECONDARY)
        _content_run(para, _split_after_colon(text))
    elif low.startswith(("practical tip:", "exam tip:")):
        style_paragraph(para)
        _label_run(para, "Tip", ACCENT_PRIMARY)
        _content_run(para, _split_after_colon(text))
    else:
        render_standard_item(para, text)

# ===============
# Bullet helpers
# ===============
def split_text_to_bullets(text: str, max_chars: int = 280) -> list[str]:
    """FIX: Increased max_chars from 240 to 280 for theory subjects."""
    if not text:
        return []
    text = text.strip()
    text = text.replace("•", "\n").replace("➢", "\n").replace(" - ", "\n- ")
    lines = []
    for part in re.split(r"\n{1,}", text):
        part = part.strip()
        if not part:
            continue
        sentences = re.split(r"(?<=[.!?])\s+", part)
        buf, cur = [], ""
        for s in sentences:
            s = s.strip()
            if not s:
                continue
            if len(cur) + len(s) + 1 <= max_chars and (len(s) < 40 or len(cur) < 40):
                cur = (cur + " " + s).strip() if cur else s
            else:
                if cur: buf.append(cur)
                cur = s
        if cur: buf.append(cur)
        lines.extend(buf)

    wrapped = []
    for line in lines:
        if len(line) <= max_chars:
            wrapped.append(line)
        else:
            words, cur = line.split(), ""
            for w in words:
                if len(cur) + len(w) + 1 <= max_chars:
                    cur = (cur + " " + w).strip()
                else:
                    wrapped.append(cur); cur = w
            if cur: wrapped.append(cur)
    return wrapped

def render_bullet_paragraph(para, text: str) -> None:
    for r in list(para.runs):
        try: r.text = ""
        except Exception: pass
    r = para.add_run()
    r.text = "➢  " + text.strip()
    r.font.name  = CONTENT_FONT_NAME
    r.font.size  = CONTENT_FONT_SIZE
    r.font.color.rgb = CONTENT_COLOR
    r.font.bold  = False
    para.level   = 0
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    try:
        para.space_before = Pt(4); para.space_after = Pt(2); para.line_spacing = 1.22
    except Exception:
        pass

# ====================
# Code block renderer
# ====================
def create_professional_code_block(slide, code_text: str,
                                   position_info: Optional[dict] = None):
    """FIX: Adaptive font sizing based on code length."""
    lines = (code_text or "").rstrip("\n").split("\n")
    line_count   = len(lines)
    max_line_len = max((len(ln) for ln in lines), default=0)

    if line_count > 30 or max_line_len > 100:
        font_pt = 11
    elif line_count > 20 or max_line_len > 80:
        font_pt = 12
    elif line_count > 12:
        font_pt = 13
    else:
        font_pt = 14

    eff_line_pt      = int(font_pt * 1.15)
    per_line_in      = eff_line_pt / 72.0
    min_height_in    = 1.2
    extra_padding_in = 0.25

    if position_info:
        left            = position_info.get("left", MARGIN_LEFT)
        top             = position_info.get("top", Inches(3))
        width           = position_info.get("width", Inches(8))
        explicit_height = position_info.get("height", None)
        max_height_in   = position_info.get("max_height_in", None)
    else:
        left            = MARGIN_LEFT + Inches(0.2)
        top             = Inches(4)
        width           = Inches(12)
        explicit_height = None
        max_height_in   = None

    text    = "\n".join(lines)
    use_dyn = explicit_height is None or (isinstance(explicit_height, str)
              and explicit_height.lower() == "auto")

    if not use_dyn:
        height = explicit_height
    else:
        shape_w_in = float(width) / float(Inches(1))
        mw = max(0.5, shape_w_in - 0.25 - 0.15)
        vl = 1
        if text:
            vl = 0
            for line in text.split("\n"):
                w = _measure_text_inches(line, MONO_FONT, float(font_pt))
                vl += max(1, int((w * 1.1 / mw) + 0.999))
        dyn = max(min_height_in, vl * per_line_in + extra_padding_in)
        if max_height_in is not None:
            dyn = min(dyn, max_height_in - 0.1)
        height = Inches(dyn)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CODE_BORDER; shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left  = Inches(0.25); tf.margin_right  = Inches(0.15)
    tf.margin_top   = Inches(0.10); tf.margin_bottom = Inches(0.10)
    tf.word_wrap    = True
    tf.auto_size    = MSO_AUTO_SIZE.NONE

    p = tf.paragraphs[0]
    p.line_spacing = Pt(eff_line_pt)
    p.space_before = p.space_after = Pt(0)
    p.alignment    = PP_PARAGRAPH_ALIGNMENT.LEFT

    r = p.add_run()
    r.text = text
    r.font.name  = MONO_FONT
    r.font.size  = Pt(font_pt)
    r.font.color.rgb = CODE_TEXT
    return shape

# ===========================
# Continuation slide factory
# ===========================
def create_continued_slide(prs, base_title: str, suffix: str = "cont.") -> SlideCtx:
    s = prs.slides.add_slide(get_blank_layout(prs))
    remove_placeholders(s)
    apply_institute_slide_background(s, prs=prs)
    add_institute_logo(s, prs=prs)
    title_txt = f"{base_title} — {suffix}" if suffix else base_title
    create_institute_title_area(s, title_txt, prs=prs, align="left")

    from .config import SAFE_SIDE_MARGIN_IN
    li = _emu_to_in(MARGIN_LEFT)
    ti = 1.00
    wi = max(2.0, _emu_to_in(prs.slide_width) - li - SAFE_SIDE_MARGIN_IN)
    bi = _emu_to_in(prs.slide_height) - 1.00

    body = s.shapes.add_textbox(Inches(li), Inches(ti), Inches(wi), Inches(bi - ti))
    tf   = body.text_frame
    tf.clear(); setup_text_frame_margins(tf)
    return SlideCtx(slide=s, tf=tf, left_in=li, top_in=ti, width_in=wi, bottom_in=bi)

def _new_cont_slide(prs, title_for_cont: Optional[str] = None):
    ns = prs.slides.add_slide(get_blank_layout(prs))
    remove_placeholders(ns)
    apply_institute_slide_background(ns, prs=prs)
    add_institute_logo(ns, prs=prs)
    if title_for_cont:
        create_institute_title_area(ns, f"{title_for_cont} (cont.)", prs=prs, align="left")
    return ns

# ==============================
# Code block pagination helper
# ==============================
def _handle_code_pagination(prs, slide_title: str, code_text: str,
                             left_in: float, top_in: float,
                             width_in: float, bottom_in: float) -> None:
    DEFAULT_FONT_PT = 14
    eff_line_pt = CODE_LINE_SPACING_PT or int(DEFAULT_FONT_PT * 1.15)
    per_line_in = eff_line_pt / 72.0
    extra_pad   = 0.25
    available_h = max(1.2, (bottom_in - top_in) - 0.15)
    mw = max(0.5, width_in - 0.25 - 0.15)
    max_vl = max(1, int((available_h - extra_pad) / per_line_in) - 2)

    def _wraps(line: str) -> int:
        w = _measure_text_inches(line, MONO_FONT, float(DEFAULT_FONT_PT))
        return max(1, int((w * 1.15 / mw) + 0.999))

    parts, bucket, used = [], [], 0
    for ln in (code_text or "").rstrip("\n").split("\n"):
        w = _wraps(ln)
        if w > max_vl:
            if bucket: parts.append("\n".join(bucket)); bucket, used = [], 0
            parts.append(ln)
        elif bucket and used + w > max_vl:
            parts.append("\n".join(bucket)); bucket, used = [ln], w
        else:
            bucket.append(ln); used += w
    if bucket: parts.append("\n".join(bucket))

    total = len(parts)
    for pi, chunk in enumerate(parts, 1):
        suffix = f"Code ({pi}/{total})" if pi > 1 else "Code"
        ctx    = create_continued_slide(prs, slide_title, suffix)
        use_dyn= chunk.count("\n") + 1 <= 12
        create_professional_code_block(ctx.slide, chunk, {
            "left": Inches(ctx.left_in), "top": Inches(ctx.top_in),
            "width": Inches(ctx.width_in),
            "height": "auto" if use_dyn else Inches(available_h - 0.1),
            "max_height_in": available_h,
        })

# ===============
# Table helpers
# ===============
def _auto_col_widths_in(header, rows, max_width_in: float,
                        font_name: str, font_size_pt: float,
                        min_col_in: float = 0.9, extra_pad_in: float = 0.30) -> list[float]:
    cols = len(header)
    natural = [0.0] * cols
    for c in range(cols):
        mx = 0.0
        for row in [header] + rows:
            txt = row[c] if c < len(row) else ""
            mx = max(mx, _measure_text_inches(txt, font_name, font_size_pt))
        natural[c] = max(min_col_in, mx + extra_pad_in)
    total = sum(natural)
    if total <= max_width_in:
        return natural
    scale = max_width_in / total
    return [max(min_col_in, w * scale) for w in natural]

def create_paginated_table(prs, slide, header, rows, col_widths,
                           left_in, top_in, max_bottom_in,
                           title_for_cont: Optional[str] = None):
    if not (prs and slide and rows and header):
        return None
    cols = len(col_widths)
    header_h, row_h = 0.5, 0.45

    def _rh(row): return max(0.35, min(0.75, len(" ".join(row)) / 100.0))

    i, cur_slide, cur_top, cur_bottom = 0, slide, top_in, max_bottom_in
    last_shape = None

    while i < len(rows):
        usable = cur_bottom - cur_top
        used, page_rows, j = header_h, [], i
        while j < len(rows):
            rh = _rh(rows[j])
            if used + rh > usable - 0.05:
                break
            page_rows.append(rows[j]); used += rh; j += 1

        if not page_rows and usable <= header_h + 0.05:
            cur_slide  = _new_cont_slide(prs, title_for_cont)
            cur_top    = 1.1
            cur_bottom = _emu_to_in(prs.slide_height) - 0.9
            continue

        tbl_shape = cur_slide.shapes.add_table(
            1 + len(page_rows), cols,
            Inches(left_in), Inches(cur_top),
            Inches(sum(col_widths)), Inches(used)
        )
        tbl = tbl_shape.table
        for c, w in enumerate(col_widths):
            tbl.columns[c].width = Inches(w)

        for c in range(cols):
            cell = tbl.cell(0, c)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = header[c] if c < len(header) else ""
            run.font.name = CONTENT_FONT_NAME; run.font.size = TITLE_FONT_SIZE
            run.font.bold = True; run.font.color.rgb = COLOR_BLACK
            cell.fill.solid(); cell.fill.fore_color.rgb = GOLD_ACCENT
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tbl.rows[0].height = Inches(header_h)

        for ri, row in enumerate(page_rows, 1):
            tbl.rows[ri].height = Inches(_rh(row))
            for c in range(cols):
                cell = tbl.cell(ri, c)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = row[c] if c < len(row) else ""
                run.font.name = CONTENT_FONT_NAME; run.font.size = CONTENT_FONT_SIZE
                run.font.color.rgb = COLOR_BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = (RGBColor(240,242,248) if ri%2==0
                                            else RGBColor(230,235,245))
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        last_shape = tbl_shape
        i += len(page_rows)
        if i < len(rows):
            cur_slide  = _new_cont_slide(prs, title_for_cont)
            cur_top    = 1.1
            cur_bottom = _emu_to_in(prs.slide_height) - 0.9

    return last_shape

def render_table_autosize(slide, prs, header, rows, left_in, top_in,
                          max_width_in, max_bottom_in, title_for_cont=None):
    fsp = float(getattr(CONTENT_FONT_SIZE, "pt", 14))
    cw  = _auto_col_widths_in(header, rows, max_width_in, CONTENT_FONT_NAME, fsp)
    return create_paginated_table(prs, slide, header, rows, cw,
                                  left_in, top_in, max_bottom_in, title_for_cont)

def render_truth_or_comparison_table(prs, base_title: str, header, rows, meta) -> None:
    is_truth = meta.get("kind") == "truth"
    if is_truth:
        table_slide_title = f"{base_title} — Truth Table" if base_title else "Truth Table"
    else:
        table_slide_title = f"{base_title} — Comparison" if base_title else "Comparison Table"
    slide = prs.slides.add_slide(get_blank_layout(prs))
    remove_placeholders(slide)
    apply_institute_slide_background(slide, prs=prs)
    add_institute_logo(slide, prs=prs)
    create_institute_title_area(slide, table_slide_title, prs=prs, align="left")

    left_in, top_in, width_in, bottom_in = _slide_box(prs)
    if not header:
        return
    cols    = len(header)
    col_w   = [max(1.0, width_in / cols)] * cols
    header_h= 0.6
    row_h   = 0.55
    max_rows_per_slide = max(4, int((bottom_in - top_in - header_h) / row_h) - 1)
    pages   = max(1, (len(rows) + max_rows_per_slide - 1) // max_rows_per_slide)

    for page in range(pages):
        pr = rows[page*max_rows_per_slide:(page+1)*max_rows_per_slide]
        s  = slide if page == 0 else _new_cont_slide(prs, base_title)
        tot= 1 + len(pr)
        ts = s.shapes.add_table(tot, cols, Inches(left_in), Inches(top_in),
                                 Inches(width_in), Inches(header_h + len(pr)*row_h))
        t  = ts.table
        for c, w in enumerate(col_w):
            t.columns[c].width = Inches(w)
        for c in range(cols):
            cell = t.cell(0, c)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = header[c]
            run.font.name = CONTENT_FONT_NAME; run.font.bold = True
            run.font.size = Pt(16); run.font.color.rgb = COLOR_BLACK
            cell.fill.solid(); cell.fill.fore_color.rgb = GOLD_ACCENT
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        t.rows[0].height = Inches(header_h)
        _is_truth = meta.get("kind") == "truth"
        for ri, row in enumerate(pr, 1):
            for c, val in enumerate(row):
                cell = t.cell(ri, c)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if _is_truth else PP_ALIGN.LEFT
                run = p.add_run(); run.text = val
                run.font.name = CONTENT_FONT_NAME; run.font.size = Pt(15)
                run.font.color.rgb = COLOR_BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = (RGBColor(245,247,252) if ri%2==0
                                            else RGBColor(230,235,245))
                cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            t.rows[ri].height = Inches(row_h)

# ==========================================
# Quiz slides — Bloom parsing + 2-per-slide
# ==========================================
def create_institute_quiz_slides(prs, quiz_items: list) -> None:
    if not quiz_items:
        return
    PER_SLIDE = 2
    blank = get_blank_layout(prs)

    for batch_start in range(0, len(quiz_items), PER_SLIDE):
        batch     = quiz_items[batch_start:batch_start + PER_SLIDE]
        slide_num = batch_start // PER_SLIDE

        slide = prs.slides.add_slide(blank)
        remove_placeholders(slide)
        apply_institute_slide_background(slide, prs=prs)
        add_institute_logo(slide, prs=prs)

        title = "Quiz" if slide_num == 0 else f"Quiz (Part {slide_num + 1})"
        create_institute_title_area(slide, title, prs=prs, align="left")

        body = slide.shapes.add_textbox(MARGIN_LEFT, Inches(1.0), Inches(12.0), Inches(5.8))
        tf   = body.text_frame
        tf.clear(); setup_text_frame_margins(tf)

        for i, q_text in enumerate(batch):
            para = (tf.paragraphs[0]
                    if (i == 0 and not tf.paragraphs[0].text)
                    else tf.add_paragraph())
            try:
                para.space_before = Pt(8)
                para.space_after  = Pt(4)
            except Exception:
                pass
            _render_quiz_question(para, q_text, batch_start + i + 1)

def _render_quiz_question(para, question_text: str, q_num: int) -> None:
    nr = para.add_run()
    nr.text = f"{q_num}. "
    nr.font.name = CONTENT_FONT_NAME; nr.font.size = Pt(18)
    nr.font.color.rgb = GOLD_ACCENT; nr.font.bold = True

    if " | " not in question_text:
        r = para.add_run()
        r.text = question_text
        r.font.name = CONTENT_FONT_NAME; r.font.size = CONTENT_FONT_SIZE
        r.font.color.rgb = CONTENT_COLOR
        return

    parts       = question_text.split(" | ")
    question    = parts[0].replace("Q:", "").strip() if parts else ""
    options     = parts[1].replace("Options:", "").strip() if len(parts) > 1 else ""
    answer      = parts[2].replace("Answer:", "").strip() if len(parts) > 2 else ""
    explanation = parts[3].replace("Explanation:", "").strip() if len(parts) > 3 else ""

    qr = para.add_run()
    qr.text = question + "\n"
    qr.font.name = CONTENT_FONT_NAME; qr.font.size = CONTENT_FONT_SIZE
    qr.font.color.rgb = CONTENT_COLOR

    if options:
        or_ = para.add_run()
        or_.text = options.replace("; ", "\n") + "\n"
        or_.font.name = CONTENT_FONT_NAME; or_.font.size = Pt(15)
        or_.font.color.rgb = CONTENT_COLOR

    if answer:
        ar = para.add_run()
        ar.text = f"Answer: {answer}"
        if explanation:
            ar.text += f"  —  {explanation}"
        ar.text += "\n"
        ar.font.name = CONTENT_FONT_NAME; ar.font.size = Pt(15)
        ar.font.color.rgb = GOLD_ACCENT

# ================
# Image insertion
# ================
def insert_image_on_slide(slide, image_path: str,
                           placement: str = "right", prs=None) -> None:
    if not image_path or not os.path.exists(image_path):
        return
    slide_w = prs.slide_width  if prs else Inches(13.33)
    slide_h = prs.slide_height if prs else Inches(7.5)

    # FIX: Adaptive image sizing
    try:
        with Image.open(image_path) as im:
            iw, ih = im.size
        aspect = ih / iw if iw else 1.0
    except Exception:
        aspect = 0.67

    img_w = Inches(4.5)
    img_h = Inches(min(3.5, 4.5 * aspect))
    top   = Inches(2.0)
    pl    = (placement or "right").lower()
    left  = (Inches(0.5) if pl == "left" else
             (slide_w - img_w) // 2 if pl == "center" else
             slide_w - img_w - Inches(0.6))
    slide.shapes.add_picture(image_path, left, top, width=img_w, height=img_h)

# ========================
# Main content renderer
# ========================
def render_enhanced_content(slide, body_shape, content_items: list, prs=None,
                             slide_title: str = "",
                             image_cache: Optional[dict] = None) -> None:

    def _extract_image_dict(item) -> Optional[dict]:
        # Case 1: already a proper dict  →  {"image": {...}}
        if isinstance(item, dict) and "image" in item:
            return item

        # Case 2: Gemini returned a JSON string instead of a dict
        #         e.g. '{"image": {"prompt": "...", "placement": "right"}}'
        if isinstance(item, str):
            s = item.strip()
            if s.startswith("{") and '"image"' in s:
                try:
                    import json
                    parsed = json.loads(s)
                    if isinstance(parsed, dict) and "image" in parsed:
                        return parsed
                except Exception:
                    pass

        return None  # not an image item

    # ── Pre-scan: detect image presence BEFORE rendering any text ─────────
    # If an image will actually be inserted, shrink the text box width NOW
    # so text never overlaps the image region.
    _IMAGE_RESERVE_IN  = 5.0
    _has_pending_image = False

    for _item in (content_items or []):
        _img_dict = _extract_image_dict(_item)
        if _img_dict:
            _prompt = _img_dict["image"].get("prompt", "")
            _path   = (image_cache or {}).get(_prompt, "")
            if _path and os.path.exists(_path):
                _has_pending_image = True
                break

    if _has_pending_image:
        _orig_width      = body_shape.width
        _new_width       = max(Inches(5.5), _orig_width - Inches(_IMAGE_RESERVE_IN))
        body_shape.width = _new_width

    # ── Frame setup ────
    tf = body_shape.text_frame
    tf.clear()
    setup_text_frame_margins(tf)

    usable_w_in    = _content_width_in(body_shape, tf)
    area_top_in    = _emu_to_in(body_shape.top)
    area_bottom_in = area_top_in + _emu_to_in(body_shape.height)

    body_font_pt = float(getattr(CONTENT_FONT_SIZE, "pt", 14))
    body_line_pt = float(BODY_LINE_SPACING_PT)
    para_gap_in  = 0.12
    used_h_in    = 0.0
    image_inserted = False

    left_in   = _emu_to_in(body_shape.left)
    top_in    = _emu_to_in(body_shape.top)
    width_in  = _emu_to_in(body_shape.width)
    bottom_in = _emu_to_in(body_shape.top + body_shape.height)

    # ── Bullet helper ────
    def _add_bullets(tf_local, bullets, used_h_local):
        nonlocal slide, tf, area_top_in, area_bottom_in, left_in, top_in, width_in, bottom_in
        for i, b in enumerate(bullets):
            para = (tf_local.paragraphs[0]
                    if (i == 0 and not tf_local.paragraphs[0].text)
                    else tf_local.add_paragraph())
            render_bullet_paragraph(para, b)
            ph = _estimate_paragraph_height_in(
                b, CONTENT_FONT_NAME, body_font_pt, usable_w_in, body_line_pt
            )
            used_h_local += ph + para_gap_in
            if (area_top_in + used_h_local) > (area_bottom_in - 0.10):
                try:
                    para.clear()
                except Exception:
                    pass
                return False, bullets[i:], used_h_local - (ph + para_gap_in)
        return True, [], used_h_local

    # ── Main render loop ────
    for idx, raw in enumerate(content_items or []):
        remaining_text: Optional[str] = None

        # ── 1. Image block ────
        # Handles both proper dict AND accidentally stringified JSON from Gemini
        img_dict = _extract_image_dict(raw)
        if img_dict is not None:
            if not image_inserted:
                img    = img_dict["image"]
                prompt = img.get("prompt", "")
                plcmt  = img.get("placement", "right")
                path   = (image_cache or {}).get(prompt)
                if path and os.path.exists(path):
                    insert_image_on_slide(slide, path, placement=plcmt, prs=prs)
                    image_inserted = True
            # Always skip — never render image JSON as text
            continue

        # ── 2. Skip non-string / empty items ────
        if not isinstance(raw, str) or not raw.strip():
            continue

        raw_stripped = raw.strip()
        raw_stripped = normalize_separators(raw_stripped)

        # ── 3. Table detection ───
        parsed_table = parse_table_from_text(raw_stripped)
        if prs and parsed_table:
            header, rows, meta = parsed_table
            kind = meta.get("kind", "table")
            if kind in {"truth", "comparison"}:
                render_truth_or_comparison_table(prs, slide_title, header, rows, meta)
            else:
                ctx = create_continued_slide(prs, slide_title or "Table", "")
                render_table_autosize(
                    slide=ctx.slide, prs=prs,
                    header=header, rows=rows,
                    left_in=ctx.left_in, top_in=ctx.top_in,
                    max_width_in=ctx.width_in, max_bottom_in=ctx.bottom_in,
                    title_for_cont=slide_title or "Table",
                )
            continue

        # ── 4. Code block ────
        code_text, remaining_text = extract_code_block(raw_stripped)
        if prs and code_text is not None:
            _handle_code_pagination(
                prs, slide_title or "Code",
                code_text, left_in, top_in, width_in, bottom_in,
            )
            if remaining_text:
                raw_stripped = remaining_text
            else:
                continue

        # ── 5. Inline pipe tables ────
        inline_tables, text_without_tables = extract_inline_pipe_tables(raw_stripped)
        if prs and inline_tables:
            for h_, r_, m_ in inline_tables:
                ctx = create_continued_slide(prs, slide_title or "Table", "")
                render_table_autosize(
                    slide=ctx.slide, prs=prs,
                    header=h_, rows=r_,
                    left_in=ctx.left_in, top_in=ctx.top_in,
                    max_width_in=ctx.width_in, max_bottom_in=ctx.bottom_in,
                    title_for_cont=slide_title or "Table",
                )
            raw_stripped = text_without_tables

        working_text = (
            remaining_text.strip()
            if isinstance(remaining_text, str) and remaining_text.strip()
            else raw_stripped
        )
        if not working_text:
            continue

        # ── 6. Semantic labels ────
        if working_text.lower().strip().startswith(_SEMANTIC_PREFIXES):
            para = (
                tf.paragraphs[0]
                if not tf.paragraphs[0].text
                else tf.add_paragraph()
            )
            render_content_item(para, working_text)
            ph = _estimate_paragraph_height_in(
                working_text, CONTENT_FONT_NAME,
                body_font_pt, usable_w_in, body_line_pt,
            )
            used_h_in += ph + para_gap_in
            continue

        # ── 7. Plain bullets ────
        bullets = split_text_to_bullets(working_text)
        ok, rem, used_h_in = _add_bullets(tf, bullets, used_h_in)
        if not ok and rem:
            while rem:
                ctx = create_continued_slide(prs, slide_title or "Topic", "cont.")
                slide, tf   = ctx.slide, ctx.tf
                area_top_in = ctx.top_in
                area_bottom_in = ctx.bottom_in
                left_in, top_in, width_in, bottom_in = (
                    ctx.left_in, ctx.top_in, ctx.width_in, ctx.bottom_in
                )
                used_h_in = 0.0
                ok2, rem2, used_h_in = _add_bullets(tf, rem, used_h_in)
                rem = rem2
