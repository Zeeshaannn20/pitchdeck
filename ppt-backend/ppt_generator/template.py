from __future__ import annotations

import os
import logging
from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

from .config import (
    TEMPLATE_PATH, SLIDE_BG_PRIMARY, GOLD_ACCENT, TITLE_BG,
    TITLE_COLOR, TITLE_FONT_NAME, TITLE_FONT_SIZE,
    CONTENT_FONT_NAME, ACCENT_PRIMARY,
    MARGIN_LEFT, SAFE_SIDE_MARGIN_IN, CONTENT_TOP_IN,
)
from .utils import _emu_to_in, _measure_text_inches

logger = logging.getLogger(__name__)

# ==================
# Template helpers
# ==================
def load_template() -> Presentation:
    if os.path.exists(TEMPLATE_PATH):
        return Presentation(TEMPLATE_PATH)
    return _create_institute_template()

def _create_institute_template() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    for layout in prs.slide_layouts:
        fill = layout.background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_PRIMARY
    return prs

def apply_institute_slide_background(slide, prs=None) -> None:
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BG_PRIMARY
        _add_gold_border(slide, prs=prs)
    except Exception as e:
        logger.warning(f"Background error: {e}")

def _add_gold_border(slide, prs=None, thickness_in: float = 0.08,
                     margin_in: float = 0.0) -> None:
    try:
        if prs is not None:
            sw, sh = prs.slide_width, prs.slide_height
        else:
            try:
                pres = slide.part.package.presentation
                sw, sh = pres.slide_width, pres.slide_height
            except Exception:
                sw, sh = Inches(13.333333), Inches(7.5)

        m  = Inches(margin_in)
        bw = Inches(thickness_in)

        outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, m, m, sw - 2*m, sh - 2*m)
        outer.fill.solid(); outer.fill.fore_color.rgb = GOLD_ACCENT; outer.line.width = Pt(0)

        iw = max(0, sw - 2*(m + bw))
        ih = max(0, sh - 2*(m + bw))
        if iw > 0 and ih > 0:
            inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, m+bw, m+bw, iw, ih)
            inner.fill.solid(); inner.fill.fore_color.rgb = SLIDE_BG_PRIMARY
            inner.line.width = Pt(0)
    except Exception as e:
        logger.warning(f"Gold border error: {e}")

def add_institute_logo(slide, position: str = "top-right", logo_path: str = "logo.png",
                       prs=None, margin_inches: float = 0.2, max_width_inches: float = 1.5,
                       max_height_inches: float = 0.7, scale: float = 1.15) -> None:
    try:
        if not os.path.exists(logo_path):
            return
        if prs is not None:
            slide_w, slide_h = prs.slide_width, prs.slide_height
        else:
            pres = slide.part.package.presentation
            slide_w, slide_h = pres.slide_width, pres.slide_height

        margin = Inches(margin_inches)
        with Image.open(logo_path) as im:
            w, h = im.size
        aspect = h / float(w) if w else 1.0

        tw = max_width_inches * scale
        th = tw * aspect
        if th > max_height_inches:
            th = max_height_inches
            tw = th / aspect

        pic = slide.shapes.add_picture(logo_path, Inches(0), Inches(0),
                                       width=Inches(tw), height=Inches(th))
        pos = (position or "top-right").lower().strip()
        if pos in ("top-right", "tr"):
            pic.left, pic.top = int(slide_w - pic.width - margin), int(margin)
        elif pos in ("top-left", "tl"):
            pic.left, pic.top = int(margin), int(margin)
        elif pos in ("bottom-right", "br"):
            pic.left, pic.top = int(slide_w - pic.width - margin), int(slide_h - pic.height - margin)
        else:
            pic.left, pic.top = int(margin), int(slide_h - pic.height - margin)
    except Exception as e:
        logger.warning(f"Logo error: {e}")

def create_institute_title_area(slide, title_text: str, subtitle_text: str = "",
                                prs=None, align: str = "left") -> None:
    try:
        slide_w  = prs.slide_width if prs is not None else Inches(13.333)
        banner_h = Inches(0.45)
        pad_lr   = Inches(0.22)
        pad_tb   = Inches(0.12)
        x_margin = Inches(0.25)

        title_clean = (title_text or "").replace("\n", " ").strip()
        tw = _measure_text_inches(title_clean, TITLE_FONT_NAME, float(TITLE_FONT_SIZE.pt))
        bw = min(Inches(tw) + 2*pad_lr + Inches(0.90), slide_w - 2*x_margin)
        bw = max(bw, Inches(3))

        a = (align or "left").lower()
        x = ((slide_w - bw) / 2 if a == "center" else
             (slide_w - bw - x_margin if a == "right" else x_margin))

        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(0.20), bw, banner_h)
        banner.fill.solid(); banner.fill.fore_color.rgb = TITLE_BG
        banner.line.color.rgb = GOLD_ACCENT; banner.line.width = Pt(2)
        banner.shadow.inherit = False

        tf = banner.text_frame
        tf.clear(); tf.word_wrap = False
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.margin_left = pad_lr; tf.margin_right = pad_lr
        tf.margin_top  = pad_tb; tf.margin_bottom = pad_tb

        p = tf.paragraphs[0]
        p.text = title_clean
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = 1.0

        r = p.runs[0]
        r.font.name = TITLE_FONT_NAME; r.font.size = TITLE_FONT_SIZE
        r.font.bold = True; r.font.color.rgb = TITLE_COLOR

        if subtitle_text:
            ps = tf.add_paragraph()
            ps.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            rs = ps.add_run()
            rs.text = subtitle_text
            rs.font.name = CONTENT_FONT_NAME; rs.font.size = Pt(16)
            rs.font.italic = True; rs.font.color.rgb = ACCENT_PRIMARY
    except Exception as e:
        logger.warning(f"Title area error: {e}")

def setup_text_frame_margins(tf) -> None:
    tf.margin_left   = Inches(0.3)
    tf.margin_right  = Inches(0.2)
    tf.margin_top    = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap     = True

def get_blank_layout(prs):
    for l in prs.slide_layouts:
        if l.name and l.name.lower() == "blank":
            return l
    for l in prs.slide_layouts:
        if len(l.placeholders) == 0:
            return l
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

def remove_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            shape.element.getparent().remove(shape.element)


