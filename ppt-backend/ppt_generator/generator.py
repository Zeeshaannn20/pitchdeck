from __future__ import annotations

import os
import logging
import tempfile
from typing import Optional

from pptx.util import Inches

from .config import MARGIN_LEFT
from .template import (
    load_template, get_blank_layout, remove_placeholders,
    apply_institute_slide_background, add_institute_logo,
    create_institute_title_area,
)
from .cover import create_ioi_cover_slide
from .renderers import render_enhanced_content, create_institute_quiz_slides
from .services import prefetch_images, create_request_temp_dir

logger = logging.getLogger(__name__)

def generate_professional_ppt(
    content: dict,
    subject: str,
    topics: str,
    special_requirements: str = "",
    instructor_name: str = "",
    qualification: str = "",
    instructor_image_b64: Optional[str] = None,
    batch_code: str = "",
) -> tuple[str, list[str]]:
    """
    Build a complete PPT from Gemini-generated content.

    Temp file isolation — every call gets its own unique temp directory:

        User A → /tmp/ppt_gen_AAA/   (images + PPT for User A)
        User B → /tmp/ppt_gen_BBB/   (images + PPT for User B)

    No two concurrent users share a directory, so cleanup of one
    request never affects another.

    The temp directory path is stored as the FIRST element of temp_files
    so that cleanup_temp_files() in routes.py removes the whole folder
    after the PPT is streamed to the client.

    Args:
        content:              Validated content dict from Gemini.
        subject:              Lecture subject string.
        topics:               Lecture topic(s) string.
        special_requirements: Any extra instructions for the deck.
        instructor_name:      Name shown on cover slide.
        qualification:        Qualification shown on cover slide.
        instructor_image_b64: Base64-encoded instructor photo (optional).
        batch_code:           Batch code shown on cover slide.

    Returns:
        (ppt_path, temp_files) where:
            ppt_path   = absolute path to the generated .pptx file
            temp_files = list of all paths to clean up after sending
                         (includes the temp dir itself as first element)
    """
    # ── Validation ──────
    if not content:
        raise ValueError("Content cannot be empty")
    if not subject:
        raise ValueError("Subject is required")
    if not topics:
        raise ValueError("Topics are required")

    request_temp_dir = create_request_temp_dir()
    temp_files: list[str] = [request_temp_dir]  # dir tracked first

    # ── Load template ─────
    prs   = load_template()
    blank = get_blank_layout(prs)

    image_cache = prefetch_images(content, temp_files, request_temp_dir)

    # ── Cover slide ──────
    create_ioi_cover_slide(
        prs,
        subject_text=subject,
        chapter_text=topics,
        instructor_name=instructor_name,
        qualification=qualification,
        instructor_image_b64=instructor_image_b64,
        logo_path="logo.png",
        batch_code=batch_code,
        temp_files=temp_files,
    )

    # ── Helper: add a standard content slide ────
    def _std_slide(title: str, items: list) -> None:
        s = prs.slides.add_slide(blank)
        remove_placeholders(s)
        apply_institute_slide_background(s, prs=prs)
        add_institute_logo(s, prs=prs)
        create_institute_title_area(s, title, prs=prs, align="left")
        body = s.shapes.add_textbox(
            MARGIN_LEFT, Inches(1.0), Inches(11), Inches(5.2)
        )
        render_enhanced_content(
            s, body, items, prs=prs,
            slide_title=title, image_cache=image_cache,
        )

    # ── Recap slide ─────
    if content.get("recap"):
        _std_slide("Previous Lecture Recap", content["recap"])

    # ── Agenda slide ────
    if content.get("agenda"):
        _std_slide("Today's Agenda", content["agenda"])

    # ── Topic slides ─────
    topic_list = content.get("topics", [])
    if len(topic_list) > 40:
        logger.warning(
            f"Large topic count ({len(topic_list)}), capping at 40 "
            f"to prevent oversized deck"
        )
        topic_list = topic_list[:40]

    for idx, topic_data in enumerate(topic_list):
        if not isinstance(topic_data, (list, tuple)) or len(topic_data) < 1:
            logger.warning(f"Skipping invalid topic at index {idx}")
            continue
        title = str(topic_data[0]).strip()
        if not title:
            logger.warning(f"Skipping empty title at index {idx}")
            continue
        _std_slide(title, list(topic_data[1:]))

    # ── Quiz slides ─────
    if content.get("quiz"):
        create_institute_quiz_slides(prs, content["quiz"])

    # ── Summary slide ────
    if content.get("summary"):
        _std_slide("Key Takeaways", content["summary"])

    # ── Guard: make sure we produced actual content ────
    if len(prs.slides) <= 1:
        raise ValueError(
            "No content slides generated — check content structure."
        )

    ppt_filename = f"output_{int(__import__('time').time() * 1000)}.pptx"
    
    ppt_path     = os.path.join(request_temp_dir, ppt_filename)
    prs.save(ppt_path)
    temp_files.append(ppt_path)

    logger.info(
        f"Generated PPT: {len(prs.slides)} slides, "
        f"{os.path.getsize(ppt_path)} bytes"
    )

    return ppt_path, temp_files
