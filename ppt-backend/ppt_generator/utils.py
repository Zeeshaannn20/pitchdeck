from __future__ import annotations

from dataclasses import dataclass
from PIL import ImageFont
from pptx.util import Inches

from .config import (
    FONT_PATH_MAP, MARGIN_LEFT,
    CONTENT_TOP_IN, CONTENT_BOTTOM_PAD, SAFE_SIDE_MARGIN_IN,
)

# =============
# Font cache
# =============
_FONT_CACHE: dict[tuple, ImageFont.FreeTypeFont] = {}

def _get_font(font_name: str, font_size_pt: float) -> ImageFont.FreeTypeFont:
    key = (font_name, int(round(font_size_pt)))
    if key not in _FONT_CACHE:
        fp = FONT_PATH_MAP.get(font_name, "arial.ttf")
        try:
            _FONT_CACHE[key] = ImageFont.truetype(fp, key[1])
        except Exception:
            _FONT_CACHE[key] = ImageFont.load_default()
    return _FONT_CACHE[key]

# ===============
# EMU helpers
# ===============
_EMU_PER_IN = int(Inches(1))

def _emu_to_in(v) -> float:
    return float(v) / float(_EMU_PER_IN)

def _in_to_emu(v: float) -> int:
    return int(Inches(v))

def _pt_to_in(pt: float) -> float:
    return pt / 72.0

# ==================
# Text measurement
# ==================
def _measure_text_inches(text: str, font_name: str, font_size_pt: float) -> float:
    font = _get_font(font_name, font_size_pt)
    try:
        bbox = font.getbbox(text or "")
        px = max(bbox[2] - bbox[0], 1)
    except Exception:
        try:
            px = max(int(font.getlength(text or "")), 1)
        except Exception:
            px = max(10 * len(text or ""), 1)
    return (px / 96.0) * 1.12

def _estimate_paragraph_height_in(text: str, font_name: str, font_size_pt: float,
                                   max_width_in: float, line_spacing_pt: float) -> float:
    t = (text or "").strip()
    if not t or max_width_in <= 0.01:
        return line_spacing_pt / 72.0
    lines = max(1, int((_measure_text_inches(t, font_name, font_size_pt) / max_width_in) + 0.999))
    return lines * (line_spacing_pt / 72.0)

def _content_width_in(shape, tf) -> float:
    w  = _emu_to_in(shape.width)
    ml = _emu_to_in(getattr(tf, "margin_left", 0))
    mr = _emu_to_in(getattr(tf, "margin_right", 0))
    return max(0.5, w - ml - mr)

def _slide_box(prs):
    sw = _emu_to_in(prs.slide_width)
    sh = _emu_to_in(prs.slide_height)
    left_in  = _emu_to_in(MARGIN_LEFT)
    top_in   = CONTENT_TOP_IN
    width_in = max(2.0, sw - left_in - SAFE_SIDE_MARGIN_IN)
    bottom_in= sh - CONTENT_BOTTOM_PAD
    return left_in, top_in, width_in, bottom_in

def _textbox(slide, left_in, top_in, width_in, height_in):
    from pptx.util import Inches
    shape = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    from .template import setup_text_frame_margins
    tf = shape.text_frame
    tf.clear()
    setup_text_frame_margins(tf)
    return shape, tf

# ====================
# SlideCtx dataclass
# ====================
@dataclass
class SlideCtx:
    slide:     object
    tf:        object
    left_in:   float
    top_in:    float
    width_in:  float
    bottom_in: float
